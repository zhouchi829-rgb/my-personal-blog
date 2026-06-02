from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色方案
COLOR_PRIMARY = RGBColor(0x43, 0x61, 0xEE)  # 主蓝色
COLOR_SECONDARY = RGBColor(0x3A, 0x0C, 0xA3)  # 深紫色
COLOR_ACCENT = RGBColor(0xF7, 0x25, 0x85)  # 粉色
COLOR_DARK = RGBColor(0x21, 0x25, 0x29)  # 深色
COLOR_LIGHT = RGBColor(0xF8, 0xF9, 0xFA)  # 浅色
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GRAY = RGBColor(0x6C, 0x75, 0x7D)
COLOR_LIGHT_GRAY = RGBColor(0xE9, 0xEC, 0xEF)
COLOR_GREEN = RGBColor(0x4C, 0xC9, 0xF0)
COLOR_ORANGE = RGBColor(0xF8, 0x96, 0x1E)
COLOR_BG = RGBColor(0xF5, 0xF7, 0xFA)

def add_background(slide, color=COLOR_BG):
    """添加背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    """添加矩形形状"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_gradient_bar(slide, left, top, width, height):
    """添加渐变条（用多个色块模拟）"""
    colors = [COLOR_PRIMARY, COLOR_SECONDARY, COLOR_ACCENT]
    bar_width = width / len(colors)
    for i, color in enumerate(colors):
        bar = add_shape(slide, left + int(bar_width * i), top, int(bar_width), height, color)
    return bar

def add_text_box(slide, left, top, width, height, text, font_size=18, color=COLOR_DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_text(slide, left, top, width, height, items, font_size=16, color=COLOR_DARK, spacing=Pt(8)):
    """添加带项目符号的文本"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = '微软雅黑'
        p.space_after = spacing
        p.level = 0
    return txBox

def add_card(slide, left, top, width, height, title, content, icon=""):
    """添加卡片样式"""
    # 卡片背景
    card = add_shape(slide, left, top, width, height, COLOR_WHITE)
    card.shadow.inherit = False
    
    # 顶部彩色条
    add_gradient_bar(slide, left, top, width, Inches(0.06))
    
    # 图标和标题
    if icon:
        add_text_box(slide, left + Inches(0.3), top + Inches(0.3), width - Inches(0.6), Inches(0.5),
                    f"{icon} {title}", font_size=18, color=COLOR_PRIMARY, bold=True)
    else:
        add_text_box(slide, left + Inches(0.3), top + Inches(0.3), width - Inches(0.6), Inches(0.5),
                    title, font_size=18, color=COLOR_PRIMARY, bold=True)
    
    # 内容
    add_text_box(slide, left + Inches(0.3), top + Inches(0.9), width - Inches(0.6), height - Inches(1.2),
                content, font_size=14, color=COLOR_GRAY)

# ========== 第1页：封面 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
add_background(slide, COLOR_WHITE)

# 顶部渐变条
add_gradient_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15))

# 左侧装饰色块
add_shape(slide, Inches(0), Inches(0), Inches(0.6), Inches(7.5), COLOR_PRIMARY)

# 主标题
add_text_box(slide, Inches(2), Inches(1.5), Inches(9), Inches(1.2),
            "个人博客网站", font_size=48, color=COLOR_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

# 副标题
add_text_box(slide, Inches(2), Inches(2.8), Inches(9), Inches(0.8),
            "—— 基于Node.js + Express + SQLite的全栈Web开发", 
            font_size=24, color=COLOR_SECONDARY, alignment=PP_ALIGN.CENTER)

# 分隔线
add_shape(slide, Inches(5), Inches(3.8), Inches(3.333), Inches(0.04), COLOR_ACCENT)

# 项目信息
add_text_box(slide, Inches(2), Inches(4.2), Inches(9), Inches(0.5),
            "毕业设计答辩", font_size=20, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(4.8), Inches(9), Inches(0.5),
            "答辩人：周驰", font_size=18, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(5.3), Inches(9), Inches(0.5),
            "完成日期：2025年12月", font_size=16, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)

# 底部装饰
add_shape(slide, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), COLOR_SECONDARY)

# ========== 第2页：目录 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_gradient_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08))

# 标题
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8),
            "目 录", font_size=36, color=COLOR_PRIMARY, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.04), COLOR_ACCENT)

# 目录项
toc_items = [
    ("01", "项目概述", "项目背景、目标与特点"),
    ("02", "技术栈", "前端、后端、数据库技术"),
    ("03", "系统架构", "整体架构与目录结构"),
    ("04", "功能模块实现", "用户认证、文章管理、图片上传、评论系统"),
    ("05", "关键技术问题", "开发中遇到的问题及解决方案"),
    ("06", "项目成果展示", "核心功能与界面展示"),
    ("07", "学习收获与总结", "技术提升与未来展望"),
]

for i, (num, title, desc) in enumerate(toc_items):
    y_pos = Inches(1.8) + Inches(0.75) * i
    
    # 编号圆圈
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y_pos, Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLOR_PRIMARY if i % 2 == 0 else COLOR_SECONDARY
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLOR_WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 标题
    add_text_box(slide, Inches(2), y_pos - Inches(0.05), Inches(4), Inches(0.4),
                title, font_size=20, color=COLOR_DARK, bold=True)
    # 描述
    add_text_box(slide, Inches(2), y_pos + Inches(0.35), Inches(6), Inches(0.3),
                desc, font_size=13, color=COLOR_GRAY)

# 右侧装饰
add_shape(slide, Inches(10), Inches(1.5), Inches(0.04), Inches(5), COLOR_LIGHT_GRAY)

# ========== 第3页：项目概述 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_gradient_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08))

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8),
            "01  项目概述", font_size=32, color=COLOR_PRIMARY, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.04), COLOR_ACCENT)

# 项目背景卡片
add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5),
        "项目背景",
        "随着互联网技术的快速发展，个人博客已成为展示个人技能、分享知识和记录生活的重要平台。本次毕业设计旨在开发一个功能完整的个人博客网站，通过实践掌握全栈Web开发技术，提升实际项目开发能力。")

# 项目目标卡片
add_card(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(2.5),
        "项目目标",
        "• 开发具有完整前后端功能的个人博客系统\n• 实现用户认证、文章管理、图片上传等核心功能\n• 构建响应式界面，提供良好的用户体验\n• 掌握Node.js、Express、SQLite等技术的实际应用")

# 项目特点卡片
add_card(slide, Inches(0.8), Inches(4.6), Inches(5.5), Inches(2.5),
        "项目特点",
        "• 技术栈全面：涵盖前端、后端、数据库全栈技术\n• 功能完整：包含用户系统、文章管理、评论系统、图片画廊等模块\n• 代码规范：遵循现代Web开发最佳实践\n• 易于部署：使用轻量级技术栈，部署简单")

# 开发周期卡片
add_card(slide, Inches(7), Inches(4.6), Inches(5.5), Inches(2.5),
        "开发周期与规模",
        "• 开发周期：2周\n• 代码行数：约3000行\n• 数据库表：4个核心表\n• API接口：15+个RESTful接口")

# ========== 第4页：技术栈 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_gradient_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08))

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8),
            "02  技术栈", font_size=32, color=COLOR_PRIMARY, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.04), COLOR_ACCENT)

# 前端技术
add_card(slide, Inches(0.8), Inches(1.8), Inches(3.7), Inches(5),
        "前端技术",
        "HTML5 / CSS3\n• 语义化标签结构\n• CSS Flexbox & Grid布局\n• 响应式设计适配多端\n\nJavaScript (ES6+)\n• Fetch API异步通信\n• DOM操作与事件处理\n• 模块化代码组织\n\nMarkdown解析\n• marked.js库\n• 实时预览功能\n• 文章目录生成")

# 后端技术
add_card(slide, Inches(4.8), Inches(1.8), Inches(3.7), Inches(5),
        "后端技术",
        "Node.js\n• JavaScript运行时环境\n• 异步非阻塞I/O\n• 丰富的npm生态\n\nExpress.js\n• 轻量级Web框架\n• RESTful API设计\n• 中间件机制\n\nJWT认证\n• JSON Web Token\n• 无状态认证\n• 令牌过期机制\n\nMulter\n• 文件上传处理\n• 文件类型验证\n• 大小限制控制")

# 数据库技术
add_card(slide, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5),
        "数据库技术",
        "SQLite\n• 嵌入式关系型数据库\n• 无需独立服务器\n• 零配置，易于部署\n\nbcryptjs\n• 密码加密存储\n• 盐值哈希算法\n• 防止彩虹表攻击\n\n数据库设计\n• 4个核心数据表\n• 外键关联约束\n• 索引优化查询\n• 事务处理支持")

# ========== 第5页：系统架构 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_gradient_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08))

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8),
            "03  系统架构", font_size=32, color=COLOR_PRIMARY, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.04), COLOR_ACCENT)

# 三层架构图
# 表示层
layer1 = add_shape(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2), COLOR_PRIMARY)
tf = layer1.text_frame
tf.paragraphs[0].text = "表示层（前端）"
tf.paragraphs[0].font.size = Pt(20)
tf.paragraphs[0].font.color.rgb = COLOR_WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

add_text_box(slide, Inches(1.5), Inches(3.1), Inches(10), Inches(0.4),
            "HTML5 / CSS3 / JavaScript  |  响应式设计  |  Fetch API  |  Markdown解析",
            font_size=13, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

# 箭头
arrow = add_shape(slide, Inches(6), Inches(3.3), Inches(1.333), Inches(0.5), COLOR_ACCENT)
tf = arrow.text_frame
tf.paragraphs[0].text = "⬇ HTTP ⬆"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = COLOR_WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# 业务逻辑层
layer2 = add_shape(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(1.2), COLOR_SECONDARY)
tf = layer2.text_frame
tf.paragraphs[0].text = "业务逻辑层（后端）"
tf.paragraphs[0].font.size = Pt(20)
tf.paragraphs[0].font.color.rgb = COLOR_WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

add_text_box(slide, Inches(1.5), Inches(5.1), Inches(10), Inches(0.4),
            "Node.js + Express.js  |  JWT认证  |  Multer文件上传  |  RESTful API",
            font_size=13, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

# 箭头
arrow2 = add_shape(slide, Inches(6), Inches(5.3), Inches(1.333), Inches(0.5), COLOR_ACCENT)
tf = arrow2.text_frame
tf.paragraphs[0].text = "⬇ SQL ⬆"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = COLOR_WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# 数据访问层
layer3 = add_shape(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(1.2), COLOR_ACCENT)
tf = layer3.text_frame
tf.paragraphs[0].text = "数据访问层（数据库）"
tf.paragraphs[0].font.size = Pt(20)
tf.paragraphs[0].font.color.rgb = COLOR_WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

add_text_box(slide, Inches(1.5), Inches(7.1), Inches(10), Inches(0.4),
            "SQLite  |  users表  |  articles表  |  comments表  |  gallery_items表",
            font_size=13, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

# ========== 第6页：功能模块 - 用户认证 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_gradient_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08))

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
            "04  功能模块实现 — 用户认证", font_size=32, color=COLOR_PRIMARY, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.04), COLOR_ACCENT)

# 功能描述
add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5),
        "功能描述",
        "• 用户注册：支持用户名、密码、邮箱注册\n• 用户登录：验证身份并返回JWT令牌\n• 密码加密：使用bcryptjs进行密码哈希存储\n• 登录状态：JWT令牌存储在localStorage\n• 安全退出：清除本地存储的令牌信息")

# 技术实现
add_card(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(2.5),
        "技术实现",
        "• JWT认证中间件验证请求令牌\n• bcryptjs密码加密（10轮盐值）\n• 令牌有效期24小时\n• 前端自动附加Authorization头\n• 路由级权限控制")

# 代码展示
code_box = add_shape(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.5), RGBColor(0x1E, 0x1E, 0x2E))
tf = code_box.text_frame
tf.word_wrap = True
code_text = """// JWT认证中间件
const authenticateToken = (req, res, next) => {
    const authHeader = req.headers['authorization'];
    const token = authHeader && authHeader.split(' ')[1];
    if (!token) return res.status(401).json({ error: '访问令牌缺失' });
    jwt.verify(token, JWT_SECRET, (err, user) => {
        if (err) return res.status(403).json({ error: '令牌无效' });
        req.user = user;
        next();
    });
};

// 用户登录
app.post('/api/login', (req, res) => {
    // 验证用户名密码 -> 生成JWT令牌 -> 返回用户信息"""
p = tf.paragraphs[0]
p.text = code_text
p.font.size = Pt(11)
p.font.color.rgb = RGBColor(0xE4, 0xE4, 0xE7)
p.font.name = 'Consolas'
p.space_after = Pt(4)

# ========== 第7页：功能模块 - 文章管理 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_gradient_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08))

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
            "04  功能模块实现 — 文章管理", font_size=32, color=COLOR_PRIMARY, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.04), COLOR_ACCENT)

# 功能描述
add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5),
        "功能描述",
        "• 文章CRUD：创建、读取、更新、删除\n• 分类管理：文章分类浏览和筛选\n• 标签系统：支持多标签标记文章\n• 搜索功能：按标题、内容、标签搜索\n• 分页展示：支持分页浏览文章列表\n• Markdown编辑：支持Markdown语法和预览")

# 数据库设计
add_card(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(2.5),
        "数据库设计",
        "• articles表包含：id, title, content, author_id, category, tags, cover_image, likes, views, status, created_at, updated_at\n• 支持文章状态管理（草稿/已发布）\n• 记录浏览量和更新时间\n• 建立作者与文章的外键关系")

# API接口列表
add_card(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.5),
        "API接口",
        "• GET /api/articles — 获取文章列表（分页）\n• GET /api/articles/:id — 获取单篇文章详情\n• POST /api/articles — 创建新文章（需认证）\n• PUT /api/articles/:id — 更新文章（需认证）\n• DELETE /api/articles/:id — 删除文章（需认证）\n• GET /api/articles/search/:query — 搜索文章\n• GET /api/categories — 获取文章分类\n• GET /api/articles/category/:category — 按分类获取文章")

# ========== 第8页：功能模块 - 图片上传与画廊 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_gradient_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08))

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
            "04  功能模块实现 — 图片上传与画廊", font_size=32, color=COLOR_PRIMARY, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.04), COLOR_ACCENT)

# 功能描述
add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5),
        "功能描述",
        "• 图片上传：支持JPG、PNG、GIF格式\n• 文件限制：最大2MB，类型验证\n• 图片预览：上传前预览和确认\n• 插入文章：上传后可直接插入文章内容\n• 封面设置：上传图片可设为文章封面\n• 画廊展示：摄影作品分类展示和筛选")

# 技术实现
add_card(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(2.5),
        "技术实现",
        "• Multer中间件处理文件上传\n• 唯一文件名生成（时间戳+随机数）\n• 文件类型和大小双重验证\n• 图片存储到public/uploads目录\n• 数据库记录图片元数据\n• 支持拖拽上传和点击上传")

# 代码展示
code_box = add_shape(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.5), RGBColor(0x1E, 0x1E, 0x2E))
tf = code_box.text_frame
tf.word_wrap = True
code_text = """// Multer文件上传配置
const upload = multer({
    storage: storage,  // 磁盘存储，唯一文件名
    limits: { fileSize: 2 * 1024 * 1024 },  // 2MB限制
    fileFilter: function (req, file, cb) {
        const allowedTypes = /jpeg|jpg|png|gif/;
        // 验证文件扩展名和MIME类型
        if (mimetype && extname) return cb(null, true);
        cb(new Error('只允许上传图片文件'));
    }
});"""
p = tf.paragraphs[0]
p.text = code_text
p.font.size = Pt(11)
p.font.color.rgb = RGBColor(0xE4, 0xE4, 0xE7)
p.font.name = 'Consolas'
p.space_after = Pt(4)

# ========== 第9页：功能模块 - 评论系统 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_gradient_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08))

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
            "04  功能模块实现 — 评论系统", font_size=32, color=COLOR_PRIMARY, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.04), COLOR_ACCENT)

# 功能描述
add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5),
        "功能描述",
        "• 匿名评论：无需登录即可发表评论\n• 评论列表：按时间倒序展示评论\n• 实时更新：发表评论后自动刷新\n• 评论计数：文章卡片显示评论数量\n• 表单验证：昵称和内容不能为空")

# 数据库设计
add_card(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(2.5),
        "数据库设计",
        "• comments表包含：id, article_id, user_name, content, created_at\n• 与articles表建立外键关联\n• 级联删除：删除文章时自动删除相关评论\n• 支持按文章ID快速查询评论")

# 交互流程
add_card(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.5),
        "交互流程",
        "1. 用户查看文章详情 → 2. 在评论表单输入昵称和内容 → 3. 点击发表评论\n4. 前端发送POST请求到 /api/articles/:id/comments → 5. 后端验证数据并存入数据库\n6. 返回成功响应 → 7. 前端刷新评论列表显示新评论\n\n特点：无需登录即可评论，降低了互动门槛，提高了用户参与度")

# ========== 第10页：关键技术问题 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_gradient_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08))

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
            "05  关键技术问题与解决方案", font_size=32, color=COLOR_PRIMARY, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.04), COLOR_ACCENT)

# 问题1
add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5),
        "问题1：图片上传后页面刷新丢失",
        "问题：上传的图片仅存储在前端JavaScript数组中，刷新后消失。\n\n解决方案：\n1. 创建gallery_items数据库表持久化存储\n2. 实现完整的图片上传API\n3. 配置multer将图片保存到服务器目录\n4. 使用唯一文件名防止冲突")

# 问题2
add_card(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(2.5),
        "问题2：用户认证令牌无效",
        "问题：上传图片时频繁出现\"令牌无效\"错误。\n\n解决方案：\n1. 改进登录流程，添加明确提示\n2. 优化令牌管理，确保正确保存到localStorage\n3. 增强错误处理，添加详细日志\n4. 在需要认证的操作前检查登录状态")

# 问题3
add_card(slide, Inches(0.8), Inches(4.6), Inches(5.5), Inches(2.5),
        "问题3：服务器无法启动",
        "问题：Node.js进程意外终止或端口被占用。\n\n解决方案：\n1. 使用正确的绝对路径启动服务器\n2. 检查端口3000是否可用\n3. 确保所有npm依赖已安装\n4. 添加详细的启动错误提示")

# 问题4
add_card(slide, Inches(7), Inches(4.6), Inches(5.5), Inches(2.5),
        "问题4：跨域访问问题",
        "问题：前端访问API时遇到跨域错误。\n\n解决方案：\n1. 在Express中配置CORS中间件\n2. 开发环境允许所有来源\n3. 生产环境限制特定域名\n4. 支持credentials跨域凭证")

# ========== 第11页：项目成果展示 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_gradient_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08))

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
            "06  项目成果展示", font_size=32, color=COLOR_PRIMARY, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.04), COLOR_ACCENT)

# 核心功能列表
features = [
    ("✅", "用户注册登录系统", "支持注册、登录、JWT认证"),
    ("✅", "文章发布和管理", "CRUD操作、分类、标签、搜索"),
    ("✅", "图片上传和画廊", "拖拽上传、分类展示、详情查看"),
    ("✅", "评论互动系统", "匿名评论、实时刷新、计数显示"),
    ("✅", "响应式网页设计", "适配桌面、平板、手机"),
    ("✅", "夜间模式切换", "一键切换、本地保存偏好"),
]

for i, (icon, title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(4) * col
    y = Inches(1.8) + Inches(2.5) * row
    
    # 卡片
    card = add_shape(slide, x, y, Inches(3.5), Inches(2.2), COLOR_WHITE)
    
    # 顶部彩色条
    colors = [COLOR_PRIMARY, COLOR_SECONDARY, COLOR_ACCENT, COLOR_GREEN, COLOR_ORANGE, COLOR_PRIMARY]
    add_shape(slide, x, y, Inches(3.5), Inches(0.06), colors[i])
    
    # 图标
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.5), Inches(0.5),
                icon, font_size=28, color=colors[i])
    
    # 标题
    add_text_box(slide, x + Inches(0.9), y + Inches(0.3), Inches(2.3), Inches(0.5),
                title, font_size=18, color=COLOR_DARK, bold=True)
    
    # 描述
    add_text_box(slide, x + Inches(0.3), y + Inches(1), Inches(2.9), Inches(1),
                desc, font_size=13, color=COLOR_GRAY)

# ========== 第12页：技术指标 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_gradient_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08))

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
            "06  技术指标与性能", font_size=32, color=COLOR_PRIMARY, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.04), COLOR_ACCENT)

# 技术指标卡片
metrics = [
    ("⏱️", "页面加载时间", "< 2秒", "优化了数据库查询和前端渲染"),
    ("👥", "并发用户支持", "50+", "异步非阻塞I/O处理请求"),
    ("⚡", "数据库响应时间", "< 100ms", "索引优化和查询优化"),
    ("📤", "图片上传速度", "< 5秒", "2MB以内图片快速上传"),
]

for i, (icon, title, value, desc) in enumerate(metrics):
    x = Inches(0.8) + Inches(3.1) * i
    card = add_shape(slide, x, Inches(1.8), Inches(2.8), Inches(2.5), COLOR_WHITE)
    add_shape(slide, x, Inches(1.8), Inches(2.8), Inches(0.06), [COLOR_PRIMARY, COLOR_SECONDARY, COLOR_ACCENT, COLOR_GREEN][i])
    
    add_text_box(slide, x + Inches(0.3), Inches(2), Inches(2.2), Inches(0.5),
                icon, font_size=32, color=[COLOR_PRIMARY, COLOR_SECONDARY, COLOR_ACCENT, COLOR_GREEN][i])
    add_text_box(slide, x + Inches(0.3), Inches(2.5), Inches(2.2), Inches(0.4),
                title, font_size=16, color=COLOR_DARK, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(2.9), Inches(2.2), Inches(0.5),
                value, font_size=28, color=COLOR_PRIMARY, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(3.5), Inches(2.2), Inches(0.5),
                desc, font_size=11, color=COLOR_GRAY)

# 性能优化措施
add_card(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.5),
        "性能优化措施",
        "数据库优化：为常用查询字段添加索引，使用连接查询减少数据库访问次数，实现分页查询避免一次性加载大量数据\n\n"
        "前端优化：压缩和合并CSS/JavaScript文件，使用图片懒加载技术，实现客户端缓存减少重复请求\n\n"
        "服务器优化：使用HTTP压缩减少传输数据量，配置适当的HTTP缓存头，实现请求限流防止恶意访问\n\n"
        "代码优化：使用异步编程避免阻塞，实现错误边界防止单点故障，添加详细日志记录便于问题排查")

# ========== 第13页：学习收获与总结 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_gradient_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08))

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
            "07  学习收获与总结", font_size=32, color=COLOR_PRIMARY, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.04), COLOR_ACCENT)

# 技术能力提升
add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5),
        "技术能力提升",
        "• 全栈开发能力：掌握了从前端到后端的完整开发流程\n• 数据库设计：学会了合理设计数据表结构和关系\n• API设计：理解了RESTful API的设计原则和实现方法\n• 安全防护：掌握了用户认证、数据验证等安全措施\n• 问题调试：提升了定位和解决技术问题的能力")

# 工程实践体会
add_card(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(2.5),
        "工程实践体会",
        "• 代码规范的重要性：良好的代码结构和注释提高了可维护性\n• 版本控制的价值：Git帮助管理代码变更和团队协作\n• 测试的必要性：充分的测试减少了线上问题的发生\n• 文档的实用性：完善的技术文档便于项目交接和维护")

# 未来展望
add_card(slide, Inches(0.8), Inches(4.6), Inches(5.5), Inches(2.5),
        "未来改进方向",
        "• 功能扩展：添加用户关注、消息通知等社交功能\n• 性能优化：引入Redis缓存、数据库读写分离\n• 移动端适配：开发移动端APP或PWA应用\n• SEO优化：改善搜索引擎友好性\n• 国际化支持：多语言界面和内容")

# 个人感悟
add_card(slide, Inches(7), Inches(4.6), Inches(5.5), Inches(2.5),
        "个人感悟",
        "通过本次毕业设计项目，我深刻体会到理论知识与实践结合的重要性。在开发过程中遇到的每一个问题都是学习的机会，每一次解决问题的过程都是能力的提升。未来，我将继续深入学习Web开发技术，不断优化和完善这个项目，同时将所学知识应用到更多的实际项目中。")

# ========== 第14页：致谢 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)

# 顶部渐变条
add_gradient_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15))

# 左侧装饰
add_shape(slide, Inches(0), Inches(0), Inches(0.6), Inches(7.5), COLOR_PRIMARY)

# 致谢内容
add_text_box(slide, Inches(2), Inches(2), Inches(9), Inches(1),
            "致  谢", font_size=48, color=COLOR_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(5.5), Inches(3.2), Inches(2.333), Inches(0.04), COLOR_ACCENT)

add_text_box(slide, Inches(2), Inches(3.6), Inches(9), Inches(1.5),
            "感谢指导老师的悉心指导\n感谢同学们的支持与帮助\n感谢各位评委老师的聆听与指导",
            font_size=22, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(5.5), Inches(9), Inches(0.5),
            "欢迎各位老师提问！", font_size=28, color=COLOR_SECONDARY, bold=True, alignment=PP_ALIGN.CENTER)

# 底部装饰
add_shape(slide, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), COLOR_SECONDARY)

# ========== 保存文件 ==========
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "毕业设计答辩PPT.pptx")
prs.save(output_path)
print(f"PPT已生成: {output_path}")
print(f"共 {len(prs.slides)} 页幻灯片")
